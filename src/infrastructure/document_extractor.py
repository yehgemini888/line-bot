"""
Infrastructure Layer: Document Text Extractor

Extracts text from various document formats (PDF, DOCX, XLSX, PPTX, CSV, TXT).
All processing is done in-memory using io.BytesIO — no temp files written to disk.
"""

import io
import csv
import logging
from dataclasses import dataclass
from typing import Optional, Dict, Callable

logger = logging.getLogger(__name__)

# Text length limit to avoid exceeding AI token limits
MAX_TEXT_LENGTH = 15000


@dataclass
class ExtractionResult:
    """Result of text extraction from a document."""
    success: bool
    text: str = ""
    page_count: Optional[int] = None
    error_message: Optional[str] = None


# Extension → MIME type mapping
EXTENSION_MIME_MAP: Dict[str, str] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc": "application/msword",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xls": "application/vnd.ms-excel",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppt": "application/vnd.ms-powerpoint",
    ".csv": "text/csv",
    ".txt": "text/plain",
}

# MIME types we can handle (modern formats only for .doc/.xls/.ppt)
SUPPORTED_MIMES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "text/csv",
    "text/plain",
}


class DocumentExtractor:
    """
    Extracts text content from document files.

    Supports: PDF, DOCX, XLSX, PPTX, CSV, TXT.
    All operations are in-memory (no disk I/O).
    """

    def is_supported(self, mime_type: str = "", filename: str = "") -> bool:
        """Check if a file type is supported for text extraction."""
        if mime_type in SUPPORTED_MIMES:
            return True
        if filename:
            detected = self.detect_mime_type(filename)
            return detected in SUPPORTED_MIMES if detected else False
        return False

    def detect_mime_type(self, filename: str) -> Optional[str]:
        """Detect MIME type from filename extension."""
        filename_lower = filename.lower()
        for ext, mime in EXTENSION_MIME_MAP.items():
            if filename_lower.endswith(ext):
                return mime
        return None

    async def extract(self, file_data: bytes, mime_type: str, ai_service=None) -> ExtractionResult:
        """
        Extract text from document bytes.

        Args:
            file_data: Raw file bytes
            mime_type: MIME type of the file
            ai_service: Optional AI service for OCR fallback (image-based PDFs)

        Returns:
            ExtractionResult with extracted text
        """
        handlers: Dict[str, Callable] = {
            "application/pdf": self._extract_pdf,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": self._extract_docx,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": self._extract_xlsx,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": self._extract_pptx,
            "text/csv": self._extract_csv,
            "text/plain": self._extract_txt,
        }

        handler = handlers.get(mime_type)
        if not handler:
            return ExtractionResult(
                success=False,
                error_message=f"不支援的檔案格式: {mime_type}"
            )

        try:
            result = handler(file_data)

            # PDF OCR fallback: if no text extracted and AI service available
            if not result.success and mime_type == "application/pdf" and ai_service:
                logger.info("🔍 [DocumentExtractor] PDF 無文字，嘗試 AI Vision OCR...")
                result = await self._extract_pdf_ocr(file_data, ai_service)

            # Truncate if too long
            if result.success and len(result.text) > MAX_TEXT_LENGTH:
                result.text = result.text[:MAX_TEXT_LENGTH] + "\n\n...(內容已截斷)"
            return result
        except Exception as e:
            return ExtractionResult(
                success=False,
                error_message=f"文字萃取失敗: {str(e)}"
            )

    def _extract_pdf(self, file_data: bytes) -> ExtractionResult:
        """Extract text from PDF."""
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(file_data))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)

        if not pages:
            return ExtractionResult(
                success=False,
                error_message="PDF 中未找到可萃取的文字（可能是掃描檔或純圖片 PDF）"
            )

        return ExtractionResult(
            success=True,
            text="\n\n".join(pages),
            page_count=len(reader.pages)
        )

    async def _extract_pdf_ocr(self, file_data: bytes, ai_service, max_pages: int = 5) -> ExtractionResult:
        """Extract text from image-based PDF using AI Vision OCR.

        Renders PDF pages as images via PyMuPDF, then sends to AI Vision.
        Limited to first `max_pages` pages to control cost/latency.
        """
        import fitz  # PyMuPDF

        try:
            doc = fitz.open(stream=file_data, filetype="pdf")
            total_pages = len(doc)
            pages_to_process = min(total_pages, max_pages)

            logger.info(f"🔍 [OCR] Processing {pages_to_process}/{total_pages} pages")

            all_text = []
            for i in range(pages_to_process):
                page = doc[i]
                # Render page as PNG image (150 DPI for balance of quality vs size)
                pix = page.get_pixmap(dpi=150)
                image_bytes = pix.tobytes("png")

                # Send to AI Vision
                result = await ai_service.analyze_image(
                    image_data=image_bytes,
                    mime_type="image/png",
                    prompt="請辨識並完整萃取這張圖片中的所有文字內容。保持原始排版格式，不要加入任何分析或摘要。"
                )

                if result.success and result.description:
                    all_text.append(f"--- 第 {i + 1} 頁 ---\n{result.description}")
                    logger.info(f"🔍 [OCR] Page {i + 1}: {len(result.description)} chars")

            doc.close()

            if not all_text:
                return ExtractionResult(
                    success=False,
                    error_message="AI Vision 也無法辨識 PDF 內容"
                )

            truncated_note = f"\n\n⚠️ 此 PDF 共 {total_pages} 頁，已透過 AI OCR 辨識前 {pages_to_process} 頁" if total_pages > max_pages else ""

            return ExtractionResult(
                success=True,
                text="\n\n".join(all_text) + truncated_note,
                page_count=total_pages
            )

        except Exception as e:
            logger.error(f"❌ [OCR] Failed: {e}")
            return ExtractionResult(
                success=False,
                error_message=f"AI Vision OCR 失敗: {str(e)}"
            )

    def _extract_docx(self, file_data: bytes) -> ExtractionResult:
        """Extract text from DOCX."""
        from docx import Document

        doc = Document(io.BytesIO(file_data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)

        if not paragraphs:
            return ExtractionResult(
                success=False,
                error_message="DOCX 中未找到文字內容"
            )

        return ExtractionResult(
            success=True,
            text="\n".join(paragraphs)
        )

    def _extract_xlsx(self, file_data: bytes) -> ExtractionResult:
        """Extract text from XLSX."""
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(file_data), read_only=True, data_only=True)
        all_text = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_rows = []
            for row in ws.iter_rows(values_only=True):
                cell_values = [str(v) for v in row if v is not None]
                if cell_values:
                    sheet_rows.append(" | ".join(cell_values))

            if sheet_rows:
                all_text.append(f"【{sheet_name}】")
                all_text.extend(sheet_rows)
                all_text.append("")  # blank line between sheets

        wb.close()

        if not all_text:
            return ExtractionResult(
                success=False,
                error_message="Excel 中未找到資料"
            )

        return ExtractionResult(
            success=True,
            text="\n".join(all_text)
        )

    def _extract_pptx(self, file_data: bytes) -> ExtractionResult:
        """Extract text from PPTX."""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_data))
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"--- 投影片 {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)
                # Also extract text from tables in slides
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_parts.append(row_text)

            if len(slide_parts) > 1:  # More than just the header
                slides_text.extend(slide_parts)

        if not slides_text:
            return ExtractionResult(
                success=False,
                error_message="簡報中未找到文字內容"
            )

        return ExtractionResult(
            success=True,
            text="\n".join(slides_text),
            page_count=len(prs.slides)
        )

    def _extract_csv(self, file_data: bytes) -> ExtractionResult:
        """Extract text from CSV."""
        # Try UTF-8 first, then fallback encodings
        text_content = None
        for encoding in ("utf-8", "utf-8-sig", "big5", "gb18030", "latin-1"):
            try:
                text_content = file_data.decode(encoding)
                break
            except (UnicodeDecodeError, ValueError):
                continue

        if text_content is None:
            return ExtractionResult(
                success=False,
                error_message="無法解碼 CSV 檔案（不支援的編碼）"
            )

        reader = csv.reader(io.StringIO(text_content))
        rows = []
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(cell.strip() for cell in row))

        if not rows:
            return ExtractionResult(
                success=False,
                error_message="CSV 中未找到資料"
            )

        return ExtractionResult(
            success=True,
            text="\n".join(rows)
        )

    def _extract_txt(self, file_data: bytes) -> ExtractionResult:
        """Extract text from plain text file."""
        for encoding in ("utf-8", "utf-8-sig", "big5", "gb18030", "latin-1"):
            try:
                text = file_data.decode(encoding)
                break
            except (UnicodeDecodeError, ValueError):
                continue
        else:
            return ExtractionResult(
                success=False,
                error_message="無法解碼文字檔案（不支援的編碼）"
            )

        text = text.strip()
        if not text:
            return ExtractionResult(
                success=False,
                error_message="文字檔案內容為空"
            )

        return ExtractionResult(
            success=True,
            text=text
        )
